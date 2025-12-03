import os
import io
import matplotlib.pyplot as plt
import seaborn as sns

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import letter
from reportlab.lib.enums import TA_LEFT


# ============================================================
# UTILITIES
# ============================================================

def disable_bullets(paragraph):
    """Remove bullet formatting in PPTX paragraphs."""
    pPr = paragraph._p.get_or_add_pPr()
    buNone = OxmlElement("a:buNone")
    pPr.append(buNone)


def clean_text(line: str):
    """Remove unwanted bullet characters but preserve numbering/headings."""
    return (
        line.replace("•", "")
            .replace("*", "")
            .replace("**", "")
            .strip()
    )


# ============================================================
# SUMMARY FORMAT
# ============================================================

def format_summary(summary_dict):
    lines = []
    lines.append(f"Total Rows: {summary_dict.get('rows')}")
    lines.append(f"Total Columns: {summary_dict.get('columns')}")
    lines.append("Missing Values:")

    for col, count in summary_dict.get("missing_values", {}).items():
        lines.append(f"{col}: {count}")

    return lines


# ============================================================
# INSIGHTS FORMAT (preserve paragraphs + numbering)
# ============================================================

def format_insights(text):
    final = []
    for line in text.split("\n"):
        line = clean_text(line)
        if len(line.strip()) > 0:
            final.append(line)
    return final


# ============================================================
# CORRELATION HEATMAP (in memory)
# ============================================================

def generate_corr_heatmap(df):
    num_df = df.select_dtypes(include="number")
    if num_df.empty:
        return None

    plt.figure(figsize=(5, 3))
    sns.heatmap(num_df.corr(), annot=True, cmap="Blues")
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight")
    buf.seek(0)
    plt.close()
    return buf


# ============================================================
# PPTX REPORT GENERATOR
# ============================================================

def create_ppt(summary, insights, plot_images, df, folder_path, filename="report.pptx"):
    prs = Presentation()
    ppt_path = os.path.join(folder_path, filename)

    # ---------------- Title Slide ----------------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Analytics Report"
    slide.placeholders[1].text = ""

    # ---------------- Summary Slide ----------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Summary"

    shape = slide.shapes.placeholders[1]
    tf = shape.text_frame
    tf.clear()

    tf.word_wrap = True
    shape.height = Inches(5)
    tf.auto_size = True

    for line in format_summary(summary):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        disable_bullets(p)

    # ---------------- Insights Slide ----------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI Insights"

    shape = slide.shapes.placeholders[1]
    tf = shape.text_frame
    tf.clear()

    tf.word_wrap = True
    shape.height = Inches(5)
    tf.auto_size = True

    for line in format_insights(insights):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        disable_bullets(p)

    # ---------------- Correlation Heatmap ----------------
    heatmap = generate_corr_heatmap(df)
    if heatmap:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Correlation Heatmap"
        slide.shapes.add_picture(heatmap, Inches(1), Inches(1.3), width=Inches(6.5))

    # ---------------- Histogram Slides ----------------
    for col, img in plot_images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Distribution: {col}"
        slide.shapes.add_picture(img, Inches(1), Inches(1.3), width=Inches(6.5))

    prs.save(ppt_path)
    return ppt_path


# ============================================================
# PDF REPORT GENERATOR
# ============================================================

def create_pdf(summary, insights, plot_images, df, folder_path, filename="report.pdf"):
    pdf_path = os.path.join(folder_path, filename)
    styles = getSampleStyleSheet()

    body_style = ParagraphStyle(
        "BodyClean",
        parent=styles["BodyText"],
        fontSize=11,
        leading=14,
        alignment=TA_LEFT,
    )

    story = []

    # ---------------- Title ----------------
    story.append(Paragraph("AI Analytics Report", styles["Title"]))
    story.append(Spacer(1, 12))

    # ---------------- Summary ----------------
    story.append(Paragraph("Dataset Summary", styles["Heading2"]))
    for line in format_summary(summary):
        story.append(Paragraph(clean_text(line), body_style))
    story.append(Spacer(1, 12))

    # ---------------- Insights ----------------
    story.append(Paragraph("AI Insights", styles["Heading2"]))
    for line in format_insights(insights):
        story.append(Paragraph(line, body_style))
    story.append(Spacer(1, 12))

    # ---------------- Heatmap ----------------
    heatmap = generate_corr_heatmap(df)
    if heatmap:
        story.append(Paragraph("Correlation Heatmap", styles["Heading3"]))
        story.append(Image(heatmap, width=350, height=250))
        story.append(Spacer(1, 12))

    # ---------------- Histograms ----------------
    for col, img in plot_images:
        story.append(Paragraph(f"Distribution: {col}", styles["Heading3"]))
        story.append(Image(img, width=350, height=250))
        story.append(Spacer(1, 12))

    pdf = SimpleDocTemplate(pdf_path, pagesize=letter)
    pdf.build(story)

    return pdf_path
